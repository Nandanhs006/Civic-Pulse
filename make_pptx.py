import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_BG = RGBColor(11, 19, 43)
    SAFFRON = RGBColor(249, 115, 22)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GREY = RGBColor(200, 214, 229)
    GREEN = RGBColor(34, 197, 94)

    def add_standard_slide(title_text, bullets):
        # Blank layout is index 6
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add background shape
        bg = slide.shapes.add_shape(
            1, # Rectangle MSO_SHAPE
            0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = SAFFRON

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        cf = content_box.text_frame
        cf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = cf.paragraphs[0]
            else:
                p = cf.add_paragraph()
            
            # Sub-bullet check
            if b.strip().startswith("- ") or b.strip().startswith("* "):
                p.level = 1
                p.text = b.strip("- *")
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.color.rgb = LIGHT_GREY
            else:
                p.level = 0
                p.text = b
                p.font.name = "Arial"
                p.font.size = Pt(20)
                p.font.color.rgb = WHITE
                p.font.bold = True
            
            p.space_after = Pt(12)

    # 1. Slide 1: Title Slide (Special layout)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "CIVIC PULSE"
    p1.font.name = "Arial"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = SAFFRON
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Citizen Ingestion Meets Real-Time PMO Governance"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf1.add_paragraph()
    p3.text = "Built on Google Cloud Run, Cloud SQL, BigQuery Federated Connections, and Gemini AI"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = LIGHT_GREY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(10)

    # Slide 2: The Problem
    add_standard_slide(
        "The Problem Statement",
        [
            "Fragmented Communication Channels",
            "- Suggestions and complaints are split across WhatsApp, email, hotlines, and local portals.",
            "High Administrative Latency (ETL Lag)",
            "- Standard analytical dashboards run batch syncs only once every 12-24 hours.",
            "Accountability Gaps",
            "- Lack of unified, real-time performance ratings for Lok Sabha MPs, Vidhan Sabha MLAs, and Ward Officers."
        ]
    )

    # Slide 3: The Solution
    add_standard_slide(
        "The Civic Pulse Solution",
        [
            "Unified Citizen Ingestion Portal",
            "- Decoupled ingest API that captures voice recordings, images, and text seamlessly.",
            "Instant AI-Routing & Triage",
            "- Classifies issues, maps GPS markers to local Wards, and flags priority instantly.",
            "Centralized PMO Command Center",
            "- Live analytics and representative directories rating resolution speed and performance."
        ]
    )

    # Slide 4: AI & Natural Language Processing
    add_standard_slide(
        "AI Pipeline: Gemini 1.5 Flash Integration",
        [
            "Dialect Voice Transcription",
            "- Transcribes regional language audio recordings into structured text formats.",
            "Grievance Category Auto-Routing",
            "- Maps reports directly to appropriate municipal departments (Water, Sanitation, Roads).",
            "Urgency Priority Weights",
            "- Computes risk scores from 1-100 to prioritize critical structural breakdowns."
        ]
    )

    # Slide 5: Decoupled Technical Architecture
    add_standard_slide(
        "Technical System Architecture",
        [
            "Stateless API Node Pool",
            "- FastAPI backend instances load balanced by Nginx and running on Google Cloud Run.",
            "Normalized Write Cache (OLTP)",
            "- Cloud SQL PostgreSQL database for concurrent ACID writes.",
            "Transactional Attachment Security",
            "- File uploads coupled directly to pre-generated Suggestion UUIDs with automatic disk/GCS cleanup if transactions fail."
        ]
    )

    # Slide 6: BigQuery Live Syncing (No ETL Lag)
    add_standard_slide(
        "BigQuery Live Federated Sync",
        [
            "Decoupling OLTP and OLAP workloads",
            "- No slow data replication, scheduling scripts, or background CRON jobs.",
            "Direct Federated Queries",
            "- BigQuery runs EXTERNAL_QUERY direct database scans to PostgreSQL in-place.",
            "Real-Time Analytics Dashboard",
            "- Access instantaneous insights on connections, saturation, and representative TAT metrics."
        ]
    )

    # Slide 7: Who It Serves
    add_standard_slide(
        "Who It Serves",
        [
            "Elected Representatives (MPs / MLAs)",
            "- Transparent oversight to monitor constituency backlog items and track public work progress.",
            "Municipal Ward Officers",
            "- Receives cleanly parsed, geolocated tasks routed straight to local administration units.",
            "Central PMO Administrators",
            "- Immediate tools to analyze regional governance trends and hold reps accountable."
        ]
    )

    # Slide 8: PMO Command Center Key Features
    add_standard_slide(
        "PMO Command Center Capabilities",
        [
            "Representative Directory",
            "- Searchable, paginated rosters listing active reps, sanction metrics, and open backlogs.",
            "Performance Index & Leaderboards",
            "- Dynamic ranking lists featuring Rank 1, 2, and 3 leaders in Gold, Blue, and Green themes.",
            "- Evaluates Governance Scores out of 100 based on resolution speed (TAT) and closure rates."
        ]
    )

    # Slide 9: Why It's Deployable Today
    add_standard_slide(
        "Why It's Deployable Today",
        [
            "Production-Grade Fail-Safes",
            "- Automated rollbacks prevent orphaned database entries and file storage leaks.",
            "Zero Migration Overhead",
            "- Integrates on top of existing local government databases without breaking downstream structures.",
            "Widescreen Responsive UI",
            "- High-fidelity glassmorphism dashboards with built-in dark modes and multi-language support."
        ]
    )

    # Slide 10: Scalability Beyond Pilot
    add_standard_slide(
        "Scalability Beyond Pilot",
        [
            "Municipal Corporation Scale (BBMP / BMC)",
            "- Establish local ward-officer networks to handle daily public sanitation and road tasks.",
            "State & National Roster Integrations",
            "- Automatically link Lok Sabha and Vidhan Sabha representatives to regional performance metrics.",
            "AI Automated Duplicate Prevention",
            "- Group identical regional reports into singular resolved actions using vector similarity search."
        ]
    )

    # Slide 11: Security & Compliance
    add_standard_slide(
        "Security & Enterprise Compliance",
        [
            "Role-Based Access Control (RBAC)",
            "- Strict user authorization gates dividing Citizen operations, MP dashboards, and PMO Admin screens.",
            "Domestic Sovereignty Support",
            "- Easily deployable on local state government servers to satisfy domestic data laws.",
            "Audit Logs",
            "- Immutable historical trail tracking status changes, officer assignments, and resolution time stamps."
        ]
    )

    prs.save("/Volumes/DiskD/Civicpulse/Civic-Pulse/Civic_Pulse_Pitch_Deck.pptx")
    print("Pitch deck PPTX generated successfully!")

if __name__ == "__main__":
    create_pitch_deck()
